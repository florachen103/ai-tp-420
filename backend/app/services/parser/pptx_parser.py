"""PPT (.pptx) 解析：按幻灯片切分，抽取标题、正文、演讲者备注。"""
from __future__ import annotations

from pptx import Presentation

from app.models.course import MaterialType
from app.services.parser.base import ParsedDocument, ParsedSection


def parse_pptx(file_path: str, filename: str) -> ParsedDocument:
    prs = Presentation(file_path)
    sections: list[ParsedSection] = []
    for i, slide in enumerate(prs.slides):
        title = f"Slide {i + 1}"
        buf: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                if shape == slide.shapes.title:
                    title = f"Slide {i + 1}: {line}"
                else:
                    buf.append(line)

        # 演讲者备注通常是讲解要点，很有价值
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                buf.append(f"\n【备注】\n{notes}")

        content = "\n".join(buf).strip()
        if content:
            sections.append(ParsedSection(
                title=title,
                content=content,
                order_idx=i,
                meta={"slide_index": i + 1, "has_notes": bool(notes)},
            ))
    return ParsedDocument(
        filename=filename,
        material_type=MaterialType.PPT,
        sections=sections,
        meta={"slide_count": len(prs.slides)},
    )
